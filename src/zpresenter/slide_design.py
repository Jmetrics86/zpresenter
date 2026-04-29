"""Slide visual system — typography scale, margins, and divider primitives."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from zpresenter.models import Deck, SlideLayoutKind


def _rgb(hex_str: str | None) -> RGBColor | None:
    if not hex_str:
        return None
    h = hex_str.strip().lstrip("#")
    if len(h) != 6:
        return None
    try:
        return RGBColor.from_string(h)
    except ValueError:
        return None


# Fallback accent when deck has no theme (professional blue).
_FALLBACK_ACCENT = RGBColor(0x25, 0x63, 0xEB)
_BODY_TEXT = RGBColor(0x1F, 0x29, 0x37)
_MUTED_DEFAULT = RGBColor(0x64, 0x74, 0x8B)


def accent_rgb(deck: Deck) -> RGBColor:
    if deck.theme and deck.theme.accent_hex:
        c = _rgb(deck.theme.accent_hex)
        if c is not None:
            return c
    return _FALLBACK_ACCENT


def muted_rgb(deck: Deck) -> RGBColor:
    if deck.theme and deck.theme.muted_hex:
        c = _rgb(deck.theme.muted_hex)
        if c is not None:
            return c
    return _MUTED_DEFAULT


def divider_line_rgb(_deck: Deck) -> RGBColor:
    """Light neutral hairline between regions."""
    return RGBColor(0xE2, 0xE8, 0xF0)


def body_text_rgb(_deck: Deck) -> RGBColor:
    return _BODY_TEXT


# --- Typography (pt) ---------------------------------------------------------

OPEN_TITLE_PT = 44
OPEN_SUBTITLE_PT = 22
SECTION_TITLE_PT = 36
CONTENT_TITLE_PT = 30
CHART_TITLE_PT = 30
CHART_SUBTITLE_PT = 15
BODY_BULLET_PT = 17
BODY_BULLET_TWO_COL_PT = 16
QUOTE_PT = 27
ATTRIBUTION_PT = 15
CAPTION_PT = 10

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

SPACE_AFTER_BULLET_PT = 10
TITLE_PARA_SPACE_BEFORE_PT = 0
TITLE_PARA_SPACE_AFTER_PT = 4


def title_pt(layout: SlideLayoutKind, *, manual_chart_box: bool = False) -> int:
    if layout == "title":
        return OPEN_TITLE_PT
    if layout == "section":
        return SECTION_TITLE_PT
    if layout == "quote":
        return QUOTE_PT
    if layout in ("chart_bar", "chart_line"):
        return CHART_TITLE_PT if manual_chart_box else CONTENT_TITLE_PT
    return CONTENT_TITLE_PT


def configure_title_textframe(tf, *, layout: SlideLayoutKind, rgb_title: RGBColor | None, deck: Deck) -> None:
    """Baseline typography on the slide title placeholder."""
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(8)
    tf.margin_bottom = Pt(4)
    para = tf.paragraphs[0]
    para.space_before = Pt(TITLE_PARA_SPACE_BEFORE_PT)
    para.space_after = Pt(TITLE_PARA_SPACE_AFTER_PT)
    para.font.name = FONT_TITLE
    para.font.bold = True
    para.font.size = Pt(title_pt(layout))
    if rgb_title:
        para.font.color.rgb = rgb_title
    else:
        para.font.color.rgb = body_text_rgb(deck)
    if layout == "quote":
        para.font.italic = True
        para.font.bold = False
    if layout == "section":
        para.alignment = PP_ALIGN.CENTER


def configure_subtitle_paragraph(paragraph, *, deck: Deck, accent_if_no_muted: bool = False) -> None:
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(OPEN_SUBTITLE_PT)
    paragraph.font.bold = False
    paragraph.space_before = Pt(4)
    paragraph.space_after = Pt(2)
    if accent_if_no_muted and deck.theme and deck.theme.accent_hex and not deck.theme.muted_hex:
        paragraph.font.color.rgb = accent_rgb(deck)
    else:
        paragraph.font.color.rgb = muted_rgb(deck)


def configure_body_paragraph(paragraph, *, deck: Deck, font_pt: int) -> None:
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(font_pt)
    paragraph.font.bold = False
    paragraph.space_after = Pt(SPACE_AFTER_BULLET_PT)
    paragraph.font.color.rgb = body_text_rgb(deck)


def configure_attribution_paragraph(paragraph, *, deck: Deck) -> None:
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(ATTRIBUTION_PT)
    paragraph.font.italic = True
    paragraph.font.color.rgb = muted_rgb(deck)
    paragraph.space_before = Pt(12)
    paragraph.alignment = PP_ALIGN.RIGHT


def apply_manual_headline(paragraph, text: str, *, rgb: RGBColor | None, deck: Deck) -> None:
    """Bold headline used on chart slides and two-column manual title boxes."""
    paragraph.text = text
    paragraph.font.name = FONT_TITLE
    paragraph.font.bold = True
    paragraph.font.italic = False
    paragraph.font.size = Pt(CONTENT_TITLE_PT)
    paragraph.space_after = Pt(10)
    paragraph.font.color.rgb = rgb if rgb else body_text_rgb(deck)


def apply_chart_subtitle_line(paragraph, text: str, *, deck: Deck) -> None:
    paragraph.text = text
    paragraph.font.name = FONT_BODY
    paragraph.font.bold = False
    paragraph.font.size = Pt(CHART_SUBTITLE_PT)
    paragraph.space_after = Pt(8)
    paragraph.font.color.rgb = muted_rgb(deck)


def configure_caption(paragraph, *, deck: Deck) -> None:
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(CAPTION_PT)
    paragraph.font.italic = True
    paragraph.font.color.rgb = muted_rgb(deck)


# --- Dividers ----------------------------------------------------------------


def add_accent_rule(s, *, left_in: float, top_in: float, width_in: float, rgb: RGBColor) -> None:
    """Short horizontal brand bar under headlines."""
    h_in = 5 / 72  # ~5 pt
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(h_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()


def add_hairline(s, *, left_in: float, top_in: float, width_in: float, rgb: RGBColor) -> None:
    """Full-width separator."""
    h_in = 1.25 / 72
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(h_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()


def add_vertical_rule(
    s,
    *,
    center_x_in: float,
    top_in: float,
    height_in: float,
    rgb: RGBColor,
) -> None:
    """Thin column divider."""
    w_in = 2 / 72
    sh = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(center_x_in - w_in / 2),
        Inches(top_in),
        Inches(w_in),
        Inches(height_in),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()


def stamp_content_header_band(
    s,
    deck: Deck,
    *,
    accent_top_in: float,
    hairline_top_in: float,
    margin_left_in: float,
    content_width_in: float,
) -> None:
    """Accent stub + hairline typical for title_content / two_column titles."""
    acc = accent_rgb(deck)
    line_rgb = divider_line_rgb(deck)
    add_accent_rule(s, left_in=margin_left_in, top_in=accent_top_in, width_in=2.95, rgb=acc)
    add_hairline(s, left_in=margin_left_in, top_in=hairline_top_in, width_in=content_width_in, rgb=line_rgb)


def stamp_section_slide_band(s, deck: Deck) -> None:
    """Centered accent stub + hairline typical for chapter markers."""
    acc_w = 3.35
    left_c = (10.0 - acc_w) / 2
    add_accent_rule(s, left_in=left_c, top_in=2.22, width_in=acc_w, rgb=accent_rgb(deck))
    add_hairline(s, left_in=MARGIN_LEFT_IN, top_in=2.31, width_in=CONTENT_WIDTH_IN, rgb=divider_line_rgb(deck))


def stamp_title_slide_footer_rule(s, deck: Deck) -> None:
    """Subtle grounding line on opening slides."""
    add_hairline(s, left_in=MARGIN_LEFT_IN, top_in=6.82, width_in=CONTENT_WIDTH_IN, rgb=divider_line_rgb(deck))


# Slide content horizontal inset (matches builder text boxes).
MARGIN_LEFT_IN = 0.54
CONTENT_WIDTH_IN = 9.08

