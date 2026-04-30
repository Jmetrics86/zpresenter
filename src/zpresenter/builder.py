"""Build PowerPoint decks from pydantic models."""

from __future__ import annotations

from collections import defaultdict
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt

from zpresenter.iconography import resolve_icon
from zpresenter.layouts_pptx import layout_for
from zpresenter.media import load_image_for_picture
from zpresenter.models import Deck, Slide, SlideLayoutKind
from zpresenter.slide_design import (
    BODY_BULLET_PT,
    BODY_BULLET_TWO_COL_PT,
    CONTENT_WIDTH_IN,
    MARGIN_LEFT_IN,
    add_vertical_rule,
    apply_chart_subtitle_line,
    apply_manual_headline,
    configure_attribution_paragraph,
    configure_body_paragraph,
    configure_caption,
    configure_subtitle_paragraph,
    configure_title_textframe,
    divider_line_rgb,
    hex_to_rgb,
    stamp_content_header_band,
    stamp_section_slide_band,
    stamp_title_slide_footer_rule,
)
from zpresenter.slide_image_layout import rects_for_images

# Narrow no-break space between catalog glyph and slide text (stable across fonts).
_ICON_GAP = "\u202f"


def _with_icon(icon_id: str | None, text: str | None) -> str:
    """Prepend catalog glyph before visible text when icon_id resolves."""
    ch = resolve_icon(icon_id)
    base = "" if text is None else text
    if not ch:
        return base
    if not base.strip():
        return ch
    return f"{ch}{_ICON_GAP}{base}"


def _decorate_lines(lines: list[str], icons: list[str | None] | None) -> list[str]:
    if not icons:
        return lines
    return [_with_icon(icons[i] if i < len(icons) else None, line) for i, line in enumerate(lines)]


_PLACEMENT_DRAW_ORDER = {
    "primary_right": 0,
    "primary_below": 1,
    "two_column_span_below": 2,
    "banner_lower": 3,
    "accent_corner": 4,
}


def _adjust_title_content_body(body, slide: Slide) -> None:
    """Resize the body placeholder to leave room for primary_right / primary_below images.

    primary_below image starts at y=5.42"; body must end by y≈5.30" to avoid overlap.
    primary_right image starts at x=5.28"; body must end by x≈5.14" to avoid overlap.
    """
    if body is None:
        return
    imgs = slide.images or []
    has_right = any(i.placement == "primary_right" for i in imgs)
    has_below = any(i.placement == "primary_below" for i in imgs)
    if has_below:
        body.top = Inches(1.42)
        body.height = Inches(3.90)   # bottom=5.32, 0.10" clear of image at 5.42
        if not has_right:
            body.left = Inches(0.52)
            body.width = Inches(9.05)
    if has_right:
        body.left = Inches(0.52)
        body.top = Inches(1.42)
        body.width = Inches(4.62)   # right=5.14, 0.14" clear of image at 5.28
        if not has_below:
            body.height = Inches(5.08)


def _inject_slide_images(s, slide: Slide, deck: Deck, asset_root: Path | None) -> None:
    imgs = slide.images or []
    if not imgs:
        return
    groups = defaultdict(list)
    for img in imgs:
        groups[img.placement].append(img)

    def sort_key(pl: str) -> int:
        return _PLACEMENT_DRAW_ORDER.get(pl, 99)

    for placement in sorted(groups.keys(), key=sort_key):
        group = groups[placement]
        rects = rects_for_images(slide.layout, placement, len(group))
        if len(rects) != len(group):
            continue
        for rect, img in zip(rects, group, strict=True):
            lo, ti, wi, he = rect
            cap = (img.caption or "").strip()
            cap_h_in = 0.26 if cap else 0.0
            pic_h = max(he - cap_h_in, 0.12)
            pic_data = load_image_for_picture(img.src, asset_root)
            s.shapes.add_picture(pic_data, Inches(lo), Inches(ti), width=Inches(wi), height=Inches(pic_h))
            if cap:
                cap_top = ti + pic_h
                tb = s.shapes.add_textbox(Inches(lo), Inches(cap_top), Inches(wi), Inches(cap_h_in))
                tb.text_frame.clear()
                cp = tb.text_frame.paragraphs[0]
                cp.text = cap.strip()
                configure_caption(cp, deck=deck)


def _rgb_for_slide_title(slide: Slide, deck: Deck, *, section_style: bool) -> RGBColor | None:
    explicit = hex_to_rgb(slide.title_color_hex)
    if explicit is not None:
        return explicit
    if section_style and deck.theme:
        return hex_to_rgb(deck.theme.primary_hex)
    return None


def _set_notes(slide, text: str | None) -> None:
    if not text:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = text


def _fill_title_placeholder(
    slide,
    title: str | None,
    rgb: RGBColor | None,
    *,
    layout_kind: SlideLayoutKind,
    deck: Deck,
) -> None:
    if slide.shapes.title is None or title is None:
        return
    tf = slide.shapes.title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    configure_title_textframe(tf, layout=layout_kind, rgb_title=rgb, deck=deck)


def _subtitle_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
            return shape
    return None


def _body_placeholder(slide):
    for shape in slide.placeholders:
        t = shape.placeholder_format.type
        if t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def _collect_body_placeholders(slide):
    bodies = []
    for shape in slide.placeholders:
        t = shape.placeholder_format.type
        if t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            bodies.append(shape)
    return bodies


def _fill_bullet_body(shape, bullets: list[str], deck: Deck, font_pt: int = BODY_BULLET_PT) -> None:
    if shape is None or not bullets:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(8)
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        configure_body_paragraph(p, deck=deck, font_pt=font_pt)


def _fill_textframe_bullets(tf, bullets: list[str], deck: Deck, font_pt: int = BODY_BULLET_TWO_COL_PT) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(8)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        configure_body_paragraph(p, deck=deck, font_pt=font_pt)


def _make_chart_data(slide: Slide) -> CategoryChartData | None:
    if not slide.chart_categories or not slide.chart_series:
        return None
    cd = CategoryChartData()
    cd.categories = slide.chart_categories
    n = len(slide.chart_categories)
    for series in slide.chart_series:
        if len(series.values) != n:
            return None
        cd.add_series(series.name, tuple(series.values))
    return cd


def _chart_type(slide: Slide) -> XL_CHART_TYPE:
    return XL_CHART_TYPE.COLUMN_CLUSTERED if slide.layout == "chart_bar" else XL_CHART_TYPE.LINE_MARKERS


def _render_chart_slide(prs: Presentation, slide: Slide, deck: Deck, asset_root: Path | None) -> None:
    layout = layout_for(slide.layout, prs)
    s = prs.slides.add_slide(layout)
    rgb = _rgb_for_slide_title(slide, deck, section_style=False)

    left_m = Inches(0.55)
    chart_top = Inches(1.42)
    chart_h = Inches(4.78) if slide.subtitle else Inches(4.88)
    chart_w = Inches(9.1)

    tb = s.shapes.add_textbox(left_m, Inches(0.38), Inches(9.05), Inches(0.92))
    tb.text_frame.clear()
    apply_manual_headline(tb.text_frame.paragraphs[0], _with_icon(slide.title_icon, slide.title or ""), rgb=rgb, deck=deck)

    if slide.subtitle:
        sub = s.shapes.add_textbox(left_m, Inches(1.08), Inches(9.05), Inches(0.44))
        sub.text_frame.clear()
        apply_chart_subtitle_line(sub.text_frame.paragraphs[0], slide.subtitle, deck=deck)
        chart_top = Inches(1.64)
        stamp_content_header_band(
            s,
            deck,
            accent_top_in=1.56,
            hairline_top_in=1.598,
            margin_left_in=MARGIN_LEFT_IN,
            content_width_in=CONTENT_WIDTH_IN,
        )
    else:
        stamp_content_header_band(
            s,
            deck,
            accent_top_in=1.14,
            hairline_top_in=1.176,
            margin_left_in=MARGIN_LEFT_IN,
            content_width_in=CONTENT_WIDTH_IN,
        )

    cd = _make_chart_data(slide)
    if cd is not None:
        s.shapes.add_chart(_chart_type(slide), left_m, chart_top, chart_w, chart_h, cd)
    _inject_slide_images(s, slide, deck, asset_root)
    _set_notes(s, slide.notes)


def _render_two_column_slide(prs: Presentation, slide: Slide, deck: Deck, asset_root: Path | None) -> None:
    layout = layout_for("two_column", prs)
    s = prs.slides.add_slide(layout)
    rgb = _rgb_for_slide_title(slide, deck, section_style=False)

    stamp_content_header_band(
        s,
        deck,
        accent_top_in=1.275,
        hairline_top_in=1.312,
        margin_left_in=MARGIN_LEFT_IN,
        content_width_in=CONTENT_WIDTH_IN,
    )

    left_bullets = _decorate_lines(slide.bullets_left or [], slide.bullets_left_icons)
    right_bullets = _decorate_lines(slide.bullets_right or [], slide.bullets_right_icons)

    span_below = any(i.placement == "two_column_span_below" for i in (slide.images or []))
    col_h_in = 4.52 if span_below else 5.05
    col_h = Inches(col_h_in)

    add_vertical_rule(
        s,
        center_x_in=5.04,
        top_in=1.38,
        height_in=col_h_in,
        rgb=divider_line_rgb(deck),
    )

    tb = s.shapes.add_textbox(Inches(0.52), Inches(0.38), Inches(9.2), Inches(0.82))
    tb.text_frame.clear()
    apply_manual_headline(tb.text_frame.paragraphs[0], _with_icon(slide.title_icon, slide.title or ""), rgb=rgb, deck=deck)

    bodies = _collect_body_placeholders(s)
    if len(bodies) >= 2:
        if span_below:
            for sh in bodies[:2]:
                sh.height = col_h
        _fill_bullet_body(bodies[0], left_bullets, deck, BODY_BULLET_TWO_COL_PT)
        _fill_bullet_body(bodies[1], right_bullets, deck, BODY_BULLET_TWO_COL_PT)
    else:
        lb = s.shapes.add_textbox(Inches(0.52), Inches(1.35), Inches(4.55), col_h)
        rb = s.shapes.add_textbox(Inches(5.22), Inches(1.35), Inches(4.55), col_h)
        _fill_textframe_bullets(lb.text_frame, left_bullets, deck, BODY_BULLET_TWO_COL_PT)
        _fill_textframe_bullets(rb.text_frame, right_bullets, deck, BODY_BULLET_TWO_COL_PT)

    _inject_slide_images(s, slide, deck, asset_root)
    _set_notes(s, slide.notes)


def _apply_slide(prs: Presentation, slide: Slide, deck: Deck, asset_root: Path | None) -> None:
    if slide.layout in ("chart_bar", "chart_line"):
        _render_chart_slide(prs, slide, deck, asset_root)
        return

    if slide.layout == "two_column":
        _render_two_column_slide(prs, slide, deck, asset_root)
        return

    layout = layout_for(slide.layout, prs)
    s = prs.slides.add_slide(layout)

    section_style = slide.layout == "section"
    rgb = _rgb_for_slide_title(slide, deck, section_style=section_style)

    if slide.layout == "title":
        _fill_title_placeholder(
            s,
            _with_icon(slide.title_icon, slide.title or ""),
            rgb,
            layout_kind="title",
            deck=deck,
        )
        sub = _subtitle_placeholder(s)
        if sub is not None and slide.subtitle:
            sub.text_frame.clear()
            sp = sub.text_frame.paragraphs[0]
            sp.text = slide.subtitle
            configure_subtitle_paragraph(sp, deck=deck, accent_if_no_muted=True)
        stamp_title_slide_footer_rule(s, deck)
        _inject_slide_images(s, slide, deck, asset_root)
        _set_notes(s, slide.notes)
        return

    if slide.layout == "section":
        _fill_title_placeholder(
            s,
            _with_icon(slide.title_icon, slide.title or ""),
            rgb,
            layout_kind="section",
            deck=deck,
        )
        stamp_section_slide_band(s, deck)
        _inject_slide_images(s, slide, deck, asset_root)
        _set_notes(s, slide.notes)
        return

    if slide.layout == "quote":
        _fill_title_placeholder(
            s,
            _with_icon(slide.title_icon, slide.quote or slide.title or ""),
            rgb,
            layout_kind="quote",
            deck=deck,
        )
        sub = _subtitle_placeholder(s)
        if sub is not None and slide.attribution:
            sub.text_frame.clear()
            ap = sub.text_frame.paragraphs[0]
            ap.text = slide.attribution
            configure_attribution_paragraph(ap, deck=deck)
        _inject_slide_images(s, slide, deck, asset_root)
        _set_notes(s, slide.notes)
        return

    # title_content
    content_rgb = _rgb_for_slide_title(slide, deck, section_style=False)
    _fill_title_placeholder(
        s,
        _with_icon(slide.title_icon, slide.title or ""),
        content_rgb,
        layout_kind="title_content",
        deck=deck,
    )
    stamp_content_header_band(
        s,
        deck,
        accent_top_in=1.29,
        hairline_top_in=1.334,
        margin_left_in=MARGIN_LEFT_IN,
        content_width_in=CONTENT_WIDTH_IN,
    )
    body = _body_placeholder(s)
    _adjust_title_content_body(body, slide)
    bullets_render = _decorate_lines(slide.bullets or [], slide.bullet_icons)
    if body is not None and bullets_render:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(8)
        for i, bullet in enumerate(bullets_render):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            configure_body_paragraph(p, deck=deck, font_pt=BODY_BULLET_PT)
    _inject_slide_images(s, slide, deck, asset_root)
    _set_notes(s, slide.notes)


def build_presentation(
    deck: Deck,
    *,
    asset_root: Path | None = None,
    template_path: Path | str | None = None,
) -> Presentation:
    if template_path is not None:
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
    props = prs.core_properties
    props.title = deck.title
    props.subject = deck.subtitle or ""
    if deck.author:
        props.author = deck.author

    for slide in deck.slides:
        _apply_slide(prs, slide, deck, asset_root)

    return prs


def save_presentation(prs: Presentation, path: Path | str) -> None:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def presentation_to_bytes(prs: Presentation) -> bytes:
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
