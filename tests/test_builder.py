import base64
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from zpresenter.builder import build_presentation, presentation_to_bytes
from zpresenter.models import Deck, Slide, SlideImage

_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _slide_visible_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def test_title_icon_prepends_glyph() -> None:
    deck = Deck(
        title="D",
        slides=[
            Slide(layout="title", title="Hello", title_icon="status.success"),
            Slide(layout="title_content", title="Agenda", bullets=["a"], bullet_icons=["data.chart"]),
        ],
    )
    raw = presentation_to_bytes(build_presentation(deck))
    prs = Presentation(BytesIO(raw))
    title_slide = prs.slides[0]
    assert title_slide.shapes.title is not None
    assert "\u2713" in title_slide.shapes.title.text
    assert "\U0001f4ca" in _slide_visible_text(prs.slides[1])


def test_slide_image_embeds_picture_shape(tmp_path: Path) -> None:
    img_path = tmp_path / "photo.png"
    img_path.write_bytes(_TINY_PNG)
    deck = Deck(
        title="T",
        slides=[
            Slide(
                layout="title",
                title="Hello",
                images=[SlideImage(src="photo.png", placement="accent_corner")],
            ),
        ],
    )
    raw = presentation_to_bytes(build_presentation(deck, asset_root=tmp_path))
    prs = Presentation(BytesIO(raw))
    slide = prs.slides[0]
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)


def test_build_roundtrip_slide_count() -> None:
    deck = Deck(
        title="Deck title",
        subtitle="Sub",
        slides=[
            Slide(layout="title", title="Hello", subtitle="World"),
            Slide(layout="title_content", title="Agenda", bullets=["One", "Two"]),
        ],
    )
    raw = presentation_to_bytes(build_presentation(deck))
    prs = Presentation(BytesIO(raw))
    assert len(prs.slides) == 2
    assert prs.core_properties.title == "Deck title"
